from django.shortcuts import render

# Create your views here.
import os
import groq
from pdfminer.high_level import extract_text
from pptx import Presentation
from django.http import JsonResponse
from rest_framework.decorators import api_view
from rest_framework.parsers import MultiPartParser, FormParser
from io import BytesIO



# Initialize Groq API
client = groq.Groq(api_key="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")  # Replace with your API key

def extract_text_from_pdf(pdf_file):
    if hasattr(pdf_file, 'read'):
        pdf_file = pdf_file.read()  # Read the file into bytes
    
    # Wrap bytes in a file-like object
    pdf_file_obj = BytesIO(pdf_file)
    
    return extract_text(pdf_file_obj).strip()

def extract_text_from_pptx(pptx_file):
    presentation = Presentation(pptx_file)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

@api_view(["POST"])
def evaluate_startup(request):
    parser_classes = (MultiPartParser, FormParser)

    pitch_deck = request.FILES.get("pitch_deck")
    team_resumes = request.FILES.getlist("team_resumes")

    if not pitch_deck:
        return JsonResponse({"error": "No pitch deck provided"}, status=400)

    # Extract text from pitch deck
    if pitch_deck.name.endswith(".pdf"):
        pitch_text = extract_text_from_pdf(pitch_deck)
    elif pitch_deck.name.endswith(".pptx"):
        pitch_text = extract_text_from_pptx(pitch_deck)
    else:
        return JsonResponse({"error": "Unsupported file format"}, status=400)

    # Extract text from resumes (PDFs only)
    resumes_text = []
    for resume in team_resumes:
        if resume.name.endswith(".pdf"):
            resumes_text.append(extract_text_from_pdf(resume))

    full_resume_text = "\n".join(resumes_text)

    # Construct evaluation prompt
    evaluation_prompt = f"""
    You are a highly experienced startup judge. Evaluate the startup based on the provided pitch deck and resumes.

    **Extracted Resume Text (Team Strength - 0-4):**
    ```
    {full_resume_text}
    ```

    **Extracted Pitch Deck Text:**
    ```
    {pitch_text}
    ```

    **Scoring Criteria (Total: 20 Points)**:
    1. **Team Strength (0-4)**
    2. **Market Opportunity (0-3)**
    3. **Innovation & Uniqueness (0-2)**
    4. **Financial Viability (0-4)**
    5. **Clarity of Idea (0-2)**
    6. **Product Viability (0-2)**
    7. **Execution Plan & Roadmap (0-2)**
    8. **Investment Readiness (0-2)**
    """

    # Call Groq API
    response = client.chat.completions.create(
        model="gemma2-9b-it",
        messages=[{"role": "system", "content": evaluation_prompt}]
    )

    evaluation_result = response.choices[0].message.content
    tokens_used = response.usage.total_tokens

    return JsonResponse({
        "evaluation": evaluation_result,
        "tokens_used": tokens_used
    })
